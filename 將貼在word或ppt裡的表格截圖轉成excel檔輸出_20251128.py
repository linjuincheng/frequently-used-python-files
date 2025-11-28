import os
import re
import tempfile
from pathlib import Path

import pandas as pd
from PIL import Image, ImageOps
import pytesseract
from pptx import Presentation
from docx import Document

# ==========================
# 使用者可調整的設定
# ==========================

# 1) PPT / Word 所在的資料夾
ROOT_DIR = r"C:\_python 2024\2024final"

# 2) Tesseract 執行檔路徑 (依你實際安裝位置修改)
pytesseract.pytesseract.tesseract_cmd = r"C:\Program Files\Tesseract-OCR\tesseract.exe"

# 3) 欄位名稱
EXPECTED_COLUMNS = [
    "MAC Address", "Vendor", "Connection",
    "Network", "WiFi", "Experience", "Signal"
]


# ==========================
# 解析一行 OCR 文字 -> 一筆資料
# ==========================

def parse_data_line(line: str):
    """
    解析一行像下面這種文字，回傳 [col1,...,col7]
    48:02:af:8a:e8:00 Aerogarden UNEO RD Default UNEO_RD Excellent -55 dBm
    """
    line = line.strip()
    if not line:
        return None
    # 表頭跳過
    if "MAC Address" in line and "Vendor" in line:
        return None

    # 清掉多餘符號
    line = line.replace("|", " ")
    if line.startswith(". "):
        line = line[2:]

    tokens = line.split()
    if not tokens:
        return None

    # 找出 MAC 位址
    mac = None
    rest = []
    for i, t in enumerate(tokens):
        if re.fullmatch(r"[0-9A-Fa-f]{2}(?::[0-9A-Fa-f]{2}){5}", t):
            mac = t
            rest = tokens[i + 1:]
            break

    if mac is None or not rest:
        return None

    # --- 從尾巴開始抓欄位 ---

    # Signal
    if rest[-1].lower() == "dbm":
        signal = rest[-2] + " " + rest[-1]
        rest = rest[:-2]
    else:
        signal = rest[-1]
        rest = rest[:-1]

    if not rest:
        return None

    # Experience
    experience = rest[-1]
    rest = rest[:-1]

    # WiFi
    wifi = rest[-1] if rest else ""
    if rest:
        rest = rest[:-1]

    # Network
    network = rest[-1] if rest else ""
    if rest:
        rest = rest[:-1]

    # 剩下的是 Vendor + Connection
    conn_idx = None
    for idx, tok in enumerate(rest):
        if tok.upper().startswith("UNEO"):  # UNEO, UNEORDPort1, ...
            conn_idx = idx
            break

    if conn_idx is None:
        vendor = " ".join(rest)
        connection = ""
    else:
        vendor = " ".join(rest[:conn_idx])
        connection = " ".join(rest[conn_idx:])

    return [mac, vendor, connection, network, wifi, experience, signal]


# ==========================
# OCR：一張圖 -> 多筆資料列
# ==========================

def ocr_table_image(image_path: Path):
    img = Image.open(image_path)

    # 灰階 + 反相，讓白底黑字比較好辨識
    gray = ImageOps.grayscale(img)
    inv = ImageOps.invert(gray)

    text = pytesseract.image_to_string(inv, lang="eng", config="--psm 6")

    rows = []
    for line in text.splitlines():
        row = parse_data_line(line)
        if row:
            rows.append(row)

    # 如果真的一筆都沒有，印出 raw OCR 方便 debug
    if not rows:
        print("  [debug] 此圖片 OCR 結果如下（沒有成功解析出表格）：")
        print(text)

    return rows


# ==========================
# 從 PPTX 抓圖片
# ==========================

def extract_images_from_pptx(pptx_path: Path, tmp_dir: Path):
    image_files = []
    prs = Presentation(pptx_path)
    idx = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.shape_type == 13:  # 13 = PICTURE
                image = shape.image
                ext = image.ext or "png"
                out_path = tmp_dir / f"{pptx_path.stem}_slide{idx}.{ext}"
                with open(out_path, "wb") as f:
                    f.write(image.blob)
                image_files.append(out_path)
                idx += 1
    return image_files


# ==========================
# 從 DOCX 抓圖片
# ==========================

def extract_images_from_docx(docx_path: Path, tmp_dir: Path):
    image_files = []
    doc = Document(docx_path)
    idx = 0
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            img_part = rel.target_part
            ext = img_part.content_type.split("/")[-1]  # image/png
            out_path = tmp_dir / f"{docx_path.stem}_img{idx}.{ext}"
            with open(out_path, "wb") as f:
                f.write(img_part.blob)
            image_files.append(out_path)
            idx += 1
    return image_files


# ==========================
# 主流程
# ==========================

def main():
    root = Path(ROOT_DIR)
    all_rows = []

    with tempfile.TemporaryDirectory() as tmp:
        tmp_dir = Path(tmp)

        # 找所有 pptx / docx
        for path in root.rglob("*"):
            if path.suffix.lower() == ".pptx":
                print(f"處理 PPTX：{path}")
                imgs = extract_images_from_pptx(path, tmp_dir)
            elif path.suffix.lower() == ".docx":
                print(f"處理 DOCX：{path}")
                imgs = extract_images_from_docx(path, tmp_dir)
            else:
                continue

            for img_path in imgs:
                print(f"  OCR 圖片：{img_path.name}")
                rows = ocr_table_image(img_path)
                all_rows.extend(rows)

    if not all_rows:
        print("沒有從任何截圖中抓到有效資料，仍請檢查截圖或 OCR 設定。")
        return

    # 建 DataFrame
    df = pd.DataFrame(all_rows, columns=EXPECTED_COLUMNS)

    # ==========================
    # 後處理：符合你的需求
    # ==========================

    # 1) MAC Address：移除冒號、改成大寫
    df["MAC Address"] = (
        df["MAC Address"]
        .astype(str)
        .str.replace(":", "", regex=False)
        .str.upper()
    )

    # 2) Signal：移除 dBm，轉成數字
    df["Signal"] = (
        df["Signal"]
        .astype(str)
        .str.replace("dBm", "", case=False, regex=False)
        .str.strip()
    )
    df["Signal"] = pd.to_numeric(df["Signal"], errors="coerce")

    # ==========================
    # 輸出 Excel
    # ==========================

    out_path = root / "wifi_merged.xlsx"
    df.to_excel(out_path, index=False)
    print(f"已輸出 Excel：{out_path}")


if __name__ == "__main__":
    main()
    
---------------------------------------------------------------------
# 1. 我有表格的截圖(如上傳檔案), 這種截圖會貼在ppt或word上, 截圖的數量可能是複數個, 這ppt或word放在C:\_python 2024\2024final 2. 要把截圖裡的表格, 我整理成excel表格,並把第1個欄位的所有小寫英文改為大寫的英文,此外,並把此欄位的:都移除掉,我要達到這種效果—(原本: 84:47:09:4c:40:46, 處理後: 8447094C4046),此外, 最後一個欄位的dBm亦都移除掉, 並把Singal欄位的儲存格的內容格式都改成數字,我要達到這種效果—(原本:-49dBm, 處理後-49),處理時,要注意欄位資料的對齊,並且不管截圖有幾個,都整合成一張excel表格輸出 3. 以上需求,請給我適當的Python程式碼
